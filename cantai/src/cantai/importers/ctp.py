"""Importador do acervo Cantai Tonos Pura (CTP)."""

import re
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation

from cantai.schemas import AuditIssue, AuditReport, Hymn

_HYMNAL = "CTP"
_DATA_TEMP_PPTX = Path(__file__).resolve().parent.parent.parent.parent / "data" / "temp" / "pptx"
_DATA_TEMP = _DATA_TEMP_PPTX.parent
_FILTROS = {"IPIVILANOVAVOTORANTIM"}
_RE_SLIDE_NUM = re.compile(r"^\d+/\d+$")
_MIN_HYMN_COUNT = 450


def _validar_pasta_origem(path: Path) -> None:
    """Valida se a pasta informada é válida para importação.

    Verifica:
    - Não é pasta temporária
    - Não contém 'temp', 'tmp', 'cache' no caminho
    - Contém arquivos PPT/PPTX suficientes
    """
    resolved = path.resolve()
    temp_resolved = _DATA_TEMP.resolve()

    # Check for temp directories
    path_lower = str(resolved).lower()
    forbidden = ["data/temp", "/temp/", "/tmp/", "/cache/", "data\\temp"]
    for f in forbidden:
        if f in path_lower:
            raise ValueError(
                f"ERRO: Pasta temporária detectada ('{f}'). "
                "Nunca use data/temp como origem. "
                "Informe a pasta ORIGINAL que contém os arquivos PPT/PPTX."
            )

    if resolved == temp_resolved or str(resolved).startswith(str(temp_resolved) + "/"):
        raise ValueError(
            "ERRO: Pasta temporária do Builder detectada. "
            "Informe a pasta ORIGINAL que contém os arquivos PPT/PPTX."
        )

    # Check file count
    ppt_count = len(list(path.glob("*.ppt")))
    pptx_count = len(list(path.rglob("*.pptx")))
    total = ppt_count + pptx_count

    if total == 0:
        raise ValueError(
            f"ERRO: Pasta '{path}' não contém arquivos PPT/PPTX."
        )

    if total < _MIN_HYMN_COUNT:
        raise ValueError(
            f"ERRO: Pasta contém apenas {total} arquivos PPT/PPTX. "
            f"Mínimo esperado: {_MIN_HYMN_COUNT}. "
            "Verifique se está apontando para a pasta CORRETA do CTP."
        )


def _convert_ppt_to_pptx(ppt_path: Path, output_dir: Path) -> bool:
    """Converte um arquivo .ppt para .pptx usando LibreOffice headless."""
    output_file = output_dir / f"{ppt_path.stem}.pptx"

    try:
        with tempfile.TemporaryDirectory() as tmp_profile:
            result = subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--norestore",
                    f"-env:UserInstallation=file://{tmp_profile}",
                    "--convert-to",
                    "pptx",
                    "--outdir",
                    str(output_dir),
                    str(ppt_path),
                ],
                capture_output=True,
                timeout=120,
                check=False,
            )
            return result.returncode == 0 and output_file.exists()
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return False


def _extrair_texto(pptx_path: Path) -> list[list[str]]:
    """Extrai texto de todos os slides de uma apresentação."""
    prs = Presentation(str(pptx_path))
    slides_texto: list[list[str]] = []

    for slide in prs.slides:
        textos_slide: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragrafo in shape.text_frame.paragraphs:
                    texto = paragrafo.text.strip().replace("\x0b", " ")
                    if texto and texto not in _FILTROS and not _RE_SLIDE_NUM.match(texto):
                        textos_slide.append(texto)
        slides_texto.append(textos_slide)

    return slides_texto


def _parse_stem(stem: str) -> tuple[str, str]:
    """Extrai número e título do stem do arquivo.

    Espera formato: '003 - DEUS, SOMENTE DEUS' -> ('003', 'DEUS, SOMENTE DEUS').
    """
    if " - " in stem:
        parts = stem.split(" - ", maxsplit=1)
        return parts[0].strip(), parts[1].strip()
    return "", stem


def _montar_hymn(pptx_path: Path, slides_texto: list[list[str]]) -> Hymn:
    """Monta um objeto Hymn a partir do caminho e texto extraído."""
    number, title = _parse_stem(pptx_path.stem)

    slides_lyrics = slides_texto[1:] if len(slides_texto) > 1 else []
    linhas = [t for slide in slides_lyrics for t in slide]
    lyrics = "\n".join(linhas)
    first_line = linhas[0] if linhas else ""

    return Hymn(
        hymnal=_HYMNAL,
        number=number,
        title=title,
        category=None,
        first_line=first_line,
        lyrics=lyrics,
        slide_count=len(slides_texto),
        source_file=pptx_path,
    )


def import_ctp(path: Path) -> tuple[list[Hymn], AuditReport]:
    """Importa acervo CTP de uma pasta.

    Retorna uma tupla (hinos, relatorio).
    """
    _validar_pasta_origem(path)

    report = AuditReport()

    _DATA_TEMP_PPTX.mkdir(parents=True, exist_ok=True)

    ppt_files = sorted(path.glob("*.ppt"))
    pptx_existentes = sorted(path.rglob("*.pptx"))

    report.ppt_encontrados = len(ppt_files)
    report.pptx_encontrados = len(pptx_existentes)

    convertidos = 0
    pptx_novos: list[Path] = []
    for ppt_file in ppt_files:
        if _convert_ppt_to_pptx(ppt_file, _DATA_TEMP_PPTX):
            convertidos += 1
            pptx_novos.append(_DATA_TEMP_PPTX / f"{ppt_file.stem}.pptx")
        else:
            report.erros.append(AuditIssue(
                arquivo=ppt_file.name,
                motivo="Falha na conversão .ppt para .pptx",
            ))
    report.convertidos = convertidos

    todos = list(pptx_novos) + list(pptx_existentes)

    vistos: set[Path] = set()
    lista_apresentacoes: list[Path] = []
    for p in todos:
        resolved = p.resolve()
        if resolved not in vistos:
            vistos.add(resolved)
            lista_apresentacoes.append(p)

    hinos: list[Hymn] = []
    for pptx_path in lista_apresentacoes:
        report.processados += 1
        stem = pptx_path.stem
        if stem.startswith("000") or "EXEMPLO" in stem.upper():
            report.ignorados.append(AuditIssue(
                arquivo=pptx_path.name,
                motivo="Arquivo de exemplo/teste",
            ))
            continue
        try:
            slides_texto = _extrair_texto(pptx_path)
            if not any(slides_texto):
                report.ignorados.append(AuditIssue(
                    arquivo=pptx_path.name,
                    motivo="Apresentação sem texto nos slides",
                ))
                continue
            hinos.append(_montar_hymn(pptx_path, slides_texto))
        except Exception as e:
            report.erros.append(AuditIssue(
                arquivo=pptx_path.name,
                motivo=f"Erro ao processar: {e}",
            ))

    report.hinos_criados = len(hinos)

    return hinos, report
